from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "outputs"
PPT_PATH = BASE_DIR / "Smart_Workforce_Presentation.pptx"


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.6), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(29, 42, 51)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(0.95), Inches(9), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(12)
        sr.font.color.rgb = RGBColor(199, 92, 42)


def add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.3), Inches(4.9))
    tf = box.text_frame
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(62, 73, 82)
        p.space_after = Pt(12)


def add_note(slide, speaker: str, note: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(1.5), Inches(5.8), Inches(4.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 252, 247)
    shape.line.color.rgb = RGBColor(225, 215, 205)
    tf = shape.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = speaker
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(199, 92, 42)
    p2 = tf.add_paragraph()
    p2.text = note
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(62, 73, 82)
    p2.space_before = Pt(10)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(120, 120, 120)


def add_chart_slide(prs: Presentation, title: str, subtitle: str, left_chart: str, right_chart: str, footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    left = OUTPUT_DIR / left_chart
    right = OUTPUT_DIR / right_chart
    if left.exists():
        slide.shapes.add_picture(str(left), Inches(0.7), Inches(1.5), width=Inches(5.7))
    if right.exists():
        slide.shapes.add_picture(str(right), Inches(6.8), Inches(1.5), width=Inches(5.4))
    add_footer(slide, footer)


def add_text_slide(prs: Presentation, title: str, bullets: list[str], speaker: str, note: str, footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_bullets(slide, bullets)
    add_note(slide, speaker, note)
    add_footer(slide, footer)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_text_slide(prs, "Smart Workforce Analytics & Management System", [
        "HCL Python Hackathon Project",
        "Database + Analytics + Charts + Flask Dashboard",
        "Presentation prepared for 3 team members",
    ], "Member 1", "Introduce the team, project title, and one-line purpose.", "Slide 1")
    add_text_slide(prs, "Problem Statement", [
        "Employee data is often scattered and difficult to manage",
        "Manual analysis takes time",
        "Management needs quick insights",
        "A smart digital system can solve this",
    ], "Member 1", "Explain the real-world problem and why your project is useful.", "Slide 2")
    add_text_slide(prs, "Project Objective", [
        "Store employee records in SQLite",
        "Perform CRUD operations",
        "Analyze data using Pandas and NumPy",
        "Generate charts with Matplotlib",
        "Present results through Flask",
    ], "Member 1", "Explain what the project was designed to do.", "Slide 3")
    add_text_slide(prs, "Technologies Used", [
        "Python",
        "SQLite",
        "Object-Oriented Programming",
        "Pandas and NumPy",
        "Matplotlib",
        "Flask",
    ], "Member 2", "Describe each technology briefly and simply.", "Slide 4")
    add_text_slide(prs, "System Architecture", [
        "Employee Data",
        "SQLite Database",
        "Python Modules",
        "Analytics Engine",
        "Charts and Dashboard",
    ], "Member 2", "Walk through the data flow from input to output.", "Slide 5")
    add_text_slide(prs, "OOP Design", [
        "Employee",
        "DatabaseManager",
        "AnalyticsEngine",
        "ReportGenerator",
    ], "Member 2", "Explain how each class has its own responsibility.", "Slide 6")
    add_text_slide(prs, "Features Implemented", [
        "Employee management",
        "Add, update, and delete operations",
        "Department salary analysis",
        "Top performer filtering",
        "Statistical calculations",
        "Web dashboard",
    ], "Member 2", "Highlight the core features completed in the project.", "Slide 7")
    add_text_slide(prs, "Dashboard Overview", [
        "Summary cards",
        "Employee form",
        "Analytics tables",
        "Charts section",
        "Management table",
    ], "Member 3", "Introduce the dashboard before the demo.", "Slide 8")
    add_chart_slide(prs, "Visual Reports", "Generated using Matplotlib", "salary_by_department.png", "salary_vs_performance.png", "Slide 9")
    add_chart_slide(prs, "More Analytics Output", "Additional charts", "salary_distribution.png", "department_workforce_share.png", "Slide 10")
    add_text_slide(prs, "What We Built", [
        "Database-driven workforce system",
        "Employee CRUD management",
        "Analytics using Pandas and NumPy",
        "Charts with Matplotlib",
        "Presentation-ready Flask app",
    ], "Member 3", "Summarize the project in a simple way.", "Slide 11")
    add_text_slide(prs, "Future Scope", [
        "Login and role-based access",
        "Search and advanced filters",
        "PDF or Excel export",
        "More HR metrics",
        "Cloud deployment",
    ], "Member 3", "Show possible improvements for the future.", "Slide 12")
    add_text_slide(prs, "Thank You", [
        "Thank you for your time",
        "We are happy to answer your questions",
    ], "All Members", "Close confidently and move to Q&A.", "Slide 13")

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    print(build_presentation())
